from pptx import Presentation
import os
import re
from moviepy import AudioFileClip
from moviepy.audio.AudioClip import concatenate_audioclips
from pptflow.utils import mylogger
import asyncio
import time

# 创建日志纪录实例
logger = mylogger.get_logger(__name__)


def tts_pause(index, length, text):
    result = text.strip()
    # 在换行符后增加停顿
    result = re.sub(r'(\r?\n)', r'\1[p1000]', result)
    # 在每一页的最后增加停顿
    result += '[p2000]'
    if index == 0:
        # 在最开头增加停顿
        result = '[p2000]' + result
    if index == length - 1:
        # 在最后增加停顿
        result += '[p10000]'
    return result


def tts_pause2(page_index, page_length, text, segment_index, segment_length):
    return text
    result = text.strip()
    # 在每一个换行符前增加停顿
    result = re.sub(r'(\r?\n)', r'[p1000]\1', result)

    # 在每一页的最后增加停顿
    if segment_index == segment_length - 1:
        result += '[p2000]'

    # 在PPT第一页和最后一页额外增加停顿
    if page_index == 0 and segment_index == 0:
        # 在ppt最开头增加停顿
        result = '[p2000]' + result
    if page_index == page_length - 1 and segment_index == segment_length - 1:
        # 在ppt最后增加停顿
        result += '[p10000]'
    return result


def tts_replace(text):
    result = text.replace('MySQL', 'My SQL[=si:kwl]') \
        .replace('负载均衡', '负载[=zai3]均衡')
    return result


def audio_file_do_replace_check(audio_file_path, note_text):
    note_text_path = audio_file_path.replace('.mp3', '.txt')
    if os.path.exists(audio_file_path) and os.path.exists(note_text_path):
        with open(note_text_path, 'r', encoding='utf-8') as note_file:
            old_note_text = note_file.read()
            if old_note_text == note_text:
                return False
    return True


def add_audio_file_text_cache(audio_file_path, note_text):
    note_text_path = audio_file_path.replace('.mp3', '.txt')
    with open(note_text_path, 'w', encoding='utf-8') as note_file:
        note_file.write(note_text)


def del_file_text_del_cache(audio_file_path):
    note_text_path = audio_file_path.replace('.mp3', '.txt')
    if os.path.exists(note_text_path):
        os.remove(note_text_path)


async def ppt_note_to_audio(tts, input_ppt_path, setting, progress_tracker=None):
    try:
        presentation = Presentation(input_ppt_path)
        file_name_without_ext = os.path.basename(input_ppt_path).split('.')[0]

        # Create a dir to save the slides as images
        if not os.path.exists(setting.audio_dir_path):
            os.makedirs(setting.audio_dir_path)

        total_slides = len(presentation.slides)
        processed_slides = 0
        # Go through each slide
        for idx, slide in enumerate(presentation.slides):
            # Checks whether the current slide falls within the specified start and end page ranges
            if setting.start_page_num and idx + 1 < setting.start_page_num:
                continue
            if setting.end_page_num and idx + 1 > setting.end_page_num:
                continue
            # Check if the slide has a notes section. If it does, to generate audio and subtitles for it.
            if slide.has_notes_slide:
                notes_slide = slide.notes_slide
                # Get the text from the notes section
                note_text = notes_slide.notes_text_frame.text
                # Generate audio and subtitles
                start_time = time.time()
                await generate_audio_and_subtitles(tts, note_text, len(presentation.slides), idx,
                                                   file_name_without_ext, setting)

                processed_slides += 1
                if progress_tracker:
                    progress = processed_slides / total_slides
                    progress_tracker.update_step(progress)

                logger.info(f"ppt_note_to_audio runtime: {time.time() - start_time:.2f} seconds")
    except Exception as e:
        logger.error(f"An error occurred: {e}", exc_info=True)
        raise e


def split_text(text, language="en", max_chars=None, max_segment_chars=None):
    # 根据语言自动设置 max_chars
    if max_chars is None:
        if language == "zh":
            max_chars = 10
            max_segment_chars = 25  # 中文每段字幕最大字符数
        elif language == "en":
            max_chars = 30
            max_segment_chars = 70  # 英文每段字幕最大字符数

        # 主要分隔符：中英文标点、换行
        delimiters = r'([,.;:!?，。；：！？\n])'

        # 按标点拆分并保留分隔符
        sentences = re.split(delimiters, text)

        # 重新组合，使分隔符紧跟前一句
        merged_sentences = []
        temp_sentence = ""

        for part in sentences:
            if re.match(delimiters, part):  # 这是标点，拼接到前面
                temp_sentence += part
            elif temp_sentence:  # 如果当前有缓存的部分，则将完整句子加入列表
                merged_sentences.append(temp_sentence.strip())
                temp_sentence = part
            else:  # 否则开始新的句子
                temp_sentence = part
        if temp_sentence:
            merged_sentences.append(temp_sentence.strip())

        # 结果分段
        result = []
        current_segment = ""

        for sentence in merged_sentences:
            if len(sentence) > max_chars:
                # 需要对过长的句子进行进一步拆分
                sub_sentences = split_long_sentence(sentence, max_chars)
                for sub in sub_sentences:
                    if len(current_segment) + len(sub) + 1 <= max_segment_chars:
                        if current_segment:
                            current_segment += " "
                        current_segment += sub
                    else:
                        result.append(current_segment)
                        current_segment = sub
            else:
                # 判断是否可以加入当前分段
                if len(current_segment) + len(sentence) + 1 <= max_segment_chars:
                    if len(current_segment) + len(sentence) + 1 <= max_chars:
                        if current_segment:
                            current_segment += " "
                        current_segment += sentence
                    else:
                        result.append(current_segment)
                        current_segment = sentence
                else:
                    result.append(current_segment)
                    current_segment = sentence

        if current_segment:
            result.append(current_segment)

        return result


def split_long_sentence(sentence, max_chars):
    """
    将过长的句子按空格或合适位置拆分，避免影响语义。
    增加了对连词、引导词和介词的特殊处理，避免拆分不必要的部分。
    优先考虑句末或句首特殊词汇进行拆分。
    """

    # 定义连词、引导词、介词的词表（可以根据需要扩展）
    conjunctions = ["and", "or", "but", "so", "yet", "for", "nor"]
    prepositions = ["by", "on", "at", "in", "with", "about", "to", "from"]
    relative_pronouns = ["who", "which", "that", "whom"]
    adverbs = ["when", "why", "where", "how", "what"]

    # 将句子按空格分割成单词
    words = sentence.split(" ")

    # 如果句子没有超过最大字符数，直接返回原句
    if len(sentence) <= max_chars:
        return [sentence]  # 如果没有超过限制，直接返回原句子

    # 如果超过最大字符限制，尝试按语义合理拆分
    result = []
    current_part = ""
    special_word_found = False  # 用于标记是否遇到特殊词汇

    # 遍历单词，尽量在连词、介词或引导词处拆分
    for word in words:
        # 如果当前部分加上单词不会超过最大字符数，继续拼接
        if len(current_part) + len(word) + 1 <= max_chars:
            if current_part:
                current_part += " "
            current_part += word
        else:
            # 如果当前部分已经达到限制，检查是否遇到特殊词汇
            if not special_word_found:
                # 如果没有找到特殊词汇，直接拼接当前部分（避免不必要拆分）
                current_part += " " + word
            else:
                # 优先考虑在连词、介词、引导词处拆分
                if any(conj in word.lower() for conj in conjunctions) or \
                        any(prep in word.lower() for prep in prepositions) or \
                        any(rel in word.lower() for rel in relative_pronouns) or \
                        any(adv in word.lower() for adv in adverbs):
                    # 如果是句末或句首的特殊词汇，合并
                    if words.index(word) == len(words) - 1 or words.index(word) == 0:
                        current_part += " " + word
                    else:
                        if current_part:
                            result.append(current_part)
                        current_part = word
                else:
                    # 如果没有找到特殊词汇，继续拼接
                    current_part += " " + word

            # 如果遇到特殊词汇，标记为已找到
            if any(conj in word.lower() for conj in conjunctions) or \
                    any(prep in word.lower() for prep in prepositions) or \
                    any(rel in word.lower() for rel in relative_pronouns) or \
                    any(adv in word.lower() for adv in adverbs):
                special_word_found = True

    # 最后处理剩余的部分
    if current_part:
        result.append(current_part)

    return result


async def generate_audio_and_subtitles(tts, text, page_length, page_index, filename_prefix, setting):
    subtitle_file, audio_clips = None, []
    text_segments = split_text(text, language=setting.language, max_segment_chars=setting.subtitle_length)
    logger.info(f'text_segments: {text_segments}')

    audio_file_path = os.path.join(setting.audio_dir_path, f"{filename_prefix}-P{page_index + 1}.mp3")
    subtitle_file_path = os.path.join(setting.audio_dir_path, f'{filename_prefix}-P{page_index + 1}.srt')

    # Check if the audio file already exists and has the same content as the current text
    if setting.audio_local_cache_enabled and not audio_file_do_replace_check(audio_file_path, ''.join(text_segments)):
        logger.info(f'{audio_file_path} already exists and the text content does not change, then skip this step.')
        return

    try:
        subtitle_file = open(subtitle_file_path, 'w')
        current_time = 0
        tasks = []  # List to hold tasks for concurrent execution

        for idx, segment_text in enumerate(text_segments):
            # file path to save the audio clip based on the current segment
            segment_audio_file_path = os.path.join(setting.audio_dir_path,
                                                   f'{filename_prefix}-P{page_index + 1}-S{idx + 1}.mp3')

            # Preprocess the text
            speech_text = tts_pause2(page_index, page_length, segment_text, idx, len(text_segments))
            speech_text = tts_replace(speech_text)

            # Add the task to the list
            tasks.append(
                generate_audio_clip(tts, speech_text, segment_audio_file_path, setting)
            )

        # Run the tasks concurrently
        results = await asyncio.gather(*tasks)

        # Process the results
        for idx, result in enumerate(results):
            if not result:
                logger.error(f"Error generating audio for segment {idx + 1}")
                continue
            segment_audio = AudioFileClip(result)
            start_time = current_time
            end_time = start_time + segment_audio.duration
            subtitle_file.write(f"{idx + 1}\n")
            subtitle_file.write(f"{format_time(start_time)} --> {format_time(end_time)}\n")
            subtitle_file.write(f"{text_segments[idx]}\n\n")
            current_time = end_time
            audio_clips.append(segment_audio)

        # Save the final audio clip to the specified path
        final_audio = concatenate_audioclips(audio_clips)
        final_audio.write_audiofile(audio_file_path, logger=None)
        add_audio_file_text_cache(audio_file_path, ''.join(text_segments))

    except Exception as e:
        logger.error(f"Error occurred: {e}", exc_info=True)
    finally:
        if subtitle_file:
            subtitle_file.close()
        for clip in audio_clips:
            clip.close()
        # Clean up temporary files
        for i in range(len(text_segments)):
            segment_audio_file_path = os.path.join(setting.audio_dir_path,
                                                   f'{filename_prefix}-P{page_index + 1}-S{i + 1}.mp3')
            if os.path.exists(segment_audio_file_path):
                os.remove(segment_audio_file_path)


async def generate_audio_clip(tts, text, output_file, setting):
    try:
        await tts(text, output_file, setting)
        return output_file  # Return the path of the generated file
    except Exception as e:
        logger.error(f"Error occurred: {e}", exc_info=True)
        return None


def format_time(seconds):
    milliseconds = int((seconds - int(seconds)) * 1000)
    minutes, seconds = divmod(int(seconds), 60)
    hours, minutes = divmod(minutes, 60)
    return f"{hours:02d}:{minutes:02d}:{seconds:02d},{milliseconds:03d}"
