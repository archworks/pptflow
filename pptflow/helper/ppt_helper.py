from pptx import Presentation

def remove_notes_from_pptx(file_path):
    # 打开PPT文件
    prs = Presentation(file_path)

    # 遍历每一张幻灯片
    for slide in prs.slides:
        # 清空备注
        slide.notes_slide.notes_text_frame.clear()

    # 保存更改到新文件
    prs.save(file_path)


def read_ppt_titles(ppt_path):
    # 打开PPT文件
    presentation = Presentation(ppt_path)
    titles = []
    # 遍历每一页
    for idx, slide in enumerate(presentation.slides):
        title = ""
        closest_top = float('inf') # 初始化为无穷大

        # 遍历所有形状
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.has_text_frame:
                # 检查文本框的顶部位置
                if shape.top < closest_top:
                    closest_top = shape.top
                    title = shape.text_frame.text # 更新标题
        print(f'P{idx+1}: {title}')
        titles.append(title)
    return titles

# 示例使用
ppt_path = r'C:\Users\kfzx-lul\Downloads\iSQL数据库灾备部署副本数对比分析.pptx'
#remove_notes_from_pptx()

titles = read_ppt_titles(ppt_path)