# Author: Valley-e
# Date: 2025/1/11  
# Description:

from pptflow.utils import mylogger, setting_dic as sd

logger = mylogger.get_logger(__name__)
logger.info("Loaded mylogger, and setting_dic")
from pptflow.utils.datapath import resource_path

import os

import json

logger.info("Loaded json")

import platform

logger.info("Loaded platform")

import shutil

logger.info("Loaded shutil")

import sys

logger.info("Loaded sys")

import re

logger.info("Loaded re")

import threading

logger.info("Loaded threading")

import customtkinter as ctk

logger.info("Loaded customtkinter")

from tkinter import PhotoImage

logger.info("Loaded PhotoImage")

from tkinter import filedialog, messagebox

logger.info("Loaded filedialog and messagebox")

from PIL import Image

logger.info("Loaded PIL.Image")

from .custom_tooltip import CustomTooltip

logger.info("Loaded CustomTooltip")

from pptflow.tts.tts_service_factory import get_tts_service

logger.info("Loaded get_tts_service")

from pptflow.config.setting_factory import get_default_setting

logger.info("Loaded get_default_setting")


class PPTFlowApp(ctk.CTk):
    def __init__(self):
        super().__init__()
        # 初始化 Setting 和 TTS
        self.setting = get_default_setting(os_name=platform.system(),
                                           tts_service_provider=os.getenv("TTS_SERVICE_PROVIDER", "kokoro").lower())
        self.tts = None
        self.current_language = 'en'
        logger.info("Current language: {}".format(self.current_language))
        self.language_modes = get_locales_subdirectories() if len(
            get_locales_subdirectories()) > 0 else sd.language_mode
        self._translations_cache = {}  # Add translation cache
        self.translations = self.get_translation()
        # 初始化 CustomTkinter
        self.theme = "light"
        ctk.set_appearance_mode(self.theme)
        ctk.set_default_color_theme("blue")

        self.step = 0
        self.labels = ["Select PPT File", "Adjust Settings", "Generate Video", "Preview and Play"]
        self.icons = ["pptfile.png", "settings.png", "generate.png", "play.png"]
        self.disabled_icons = ["pptfile-disabled.png", "settings-disabled.png", "generate-disabled.png",
                               "play-disabled.png"]
        self.arrow_icons = ["arrow.png", "arrow-disabled.png"]
        # Get icon directory path
        self.icon_dir = resource_path(os.path.join("assets", "icons"))
        self.iconphoto(False, PhotoImage(file=os.path.join(self.icon_dir, "pptflow.png")))

        self.loading_title = self.get_text("generate_video")
        self.file_display = ""
        self.tooltip = None
        self.cancel_file = None
        self.select_button = None
        self.adjust_settings = None
        self.cancel_settings = None
        self.generate_button = None

        # Configure window
        self.title("PPTFlow")
        self.center_window()

        self.grid_columnconfigure(0, weight=1)
        self.grid_rowconfigure(0, weight=1)

        # Create main frame
        self.main_frame = ctk.CTkFrame(self, border_width=0, corner_radius=0, fg_color="transparent")
        self.main_frame.grid(row=0, column=0, sticky="nsew")
        self.main_frame.grid_columnconfigure(0, weight=1)
        self.main_frame.grid_rowconfigure(0, weight=1)
        self.main_frame.grid_rowconfigure(1, weight=3)
        self.main_frame.grid_rowconfigure(2, weight=1)

        # Load and set background image
        bg_image = ctk.CTkImage(Image.open(os.path.join(self.icon_dir, "background.png")), size=(800, 173))
        self.bg_label = ctk.CTkLabel(self.main_frame, image=bg_image, text="")
        self.bg_label.place(x=0, y=0, anchor="nw")

        # 顶部标题和图标部分
        self.create_top_section()

        # 中间流程图布局
        self.create_workflow_section()
        logger.info('initiated PPTFlowApp')

    def create_top_section(self):
        # 右上角的图标按钮（设置和GitHub）
        self.settings_button = ctk.CTkButton(self.main_frame, text="",
                                             image=self.load_ctk_image(
                                                 os.path.join(self.icon_dir, "system-setting.png"),
                                                 size=20),
                                             width=20, height=20, fg_color="#0066FF", corner_radius=0,
                                             hover=False,
                                             command=lambda: self.select_frame("System Settings"))
        self.help_button = ctk.CTkButton(self.main_frame, text="",
                                         image=self.load_ctk_image(os.path.join(self.icon_dir, "help.png"),
                                                                   size=20),
                                         width=20, height=20, fg_color="#0066FF", corner_radius=0,
                                         hover=False,
                                         command=lambda: self.on_button_click("Website"))

        self.settings_button.grid(row=0, column=1, padx=0, pady=10, sticky="ne")
        self.help_button.grid(row=0, column=2, padx=(0, 5), pady=10, sticky="ne")

    def create_workflow_section(self):
        self.flow_frame = ctk.CTkFrame(self.main_frame, corner_radius=0, fg_color="white")
        self.flow_frame.grid(row=1, column=0, columnspan=7, padx=100, pady=(40, 60), sticky="nsew")
        for i in range(7):
            self.flow_frame.grid_columnconfigure(i, weight=1)
        self.icon_frame = ctk.CTkFrame(self.flow_frame, corner_radius=0, fg_color="transparent")
        self.icon_frame.grid(row=0, column=0, columnspan=7, sticky="nsew")
        for i in range(7):
            self.icon_frame.grid_columnconfigure(i, weight=1)
        self.button_frame = ctk.CTkFrame(self.flow_frame, corner_radius=0, border_width=0, fg_color="transparent")
        self.button_frame.grid(row=1, column=0, columnspan=7, sticky="nsew")

        # 流程图布局
        self.select_flow_0(0)
        self.setting_flow_1(1)
        self.generation_flow_2(2)
        self.review_flow_3(3)

    def select_flow_0(self, i, row_offset=1):
        icon_image, arrow_image = self.update_image(i)
        pptfile_icon = ctk.CTkButton(self.icon_frame, image=icon_image, text="", width=40, height=60,
                                     fg_color="transparent", hover=False)
        pptfile_icon.grid(row=0, column=i * 2, padx=(60, 0), pady=(50, 0))
        arrow = ctk.CTkButton(self.icon_frame, text='', width=50, height=20, image=arrow_image,
                              fg_color="transparent", hover=False)
        arrow.grid(row=0, column=i * 2 + 1, pady=(50, 0))
        self.file_label_var = ctk.StringVar()
        self.file_label = ctk.CTkLabel(self.button_frame, textvariable=self.file_label_var,
                                       font=ctk.CTkFont(size=12), width=100)
        self.file_label.grid(row=row_offset + 1, column=i * 2, padx=(40, 5), pady=5)

        self.cancel_file = ctk.CTkButton(self.button_frame, text=self.get_text("cancel"),
                                         font=ctk.CTkFont(size=12), width=110, border_width=1, hover=False,
                                         fg_color="transparent", text_color="#2563EB", border_color="#2563EB",
                                         command=lambda: self.on_button_click("Cancel File"))
        self.cancel_file.grid(row=row_offset + 2, column=i * 2, padx=(40, 5), pady=5)
        self.file_label.grid_remove()
        self.cancel_file.grid_remove()

        self.select_button = ctk.CTkButton(self.button_frame, text=self.get_text("select_ppt"),
                                           font=ctk.CTkFont(size=12), width=110, hover_color="#1D4ED8",
                                           command=lambda: self.on_button_click("Select PPT File"))
        self.select_button.grid(row=row_offset + 1, column=i * 2, padx=(40, 5), pady=5)
        self.update_button(i, self.select_button)

    def setting_flow_1(self, i, row_offset=1):
        icon_image, arrow_image = self.update_image(i)
        setting_icon = ctk.CTkButton(self.icon_frame, image=icon_image, text="", width=65,
                                     height=65, fg_color="transparent", hover=False)
        setting_icon.grid(row=0, column=i * 2, padx=0, pady=(50, 0))
        arrow = ctk.CTkButton(self.icon_frame, text='', width=50, height=20, image=arrow_image,
                              fg_color="transparent", hover=False)
        arrow.grid(row=0, column=i * 2 + 1, pady=(50, 0))
        self.adjust_button = ctk.CTkButton(self.button_frame, text=self.get_text("adjust_settings"),
                                           font=ctk.CTkFont(size=12), width=110,
                                           text_color="white", text_color_disabled="white",
                                           command=lambda: self.select_frame("Adjust Settings"))
        self.skip_button = ctk.CTkButton(self.button_frame, text=self.get_text("skip_settings"),
                                         font=ctk.CTkFont(size=12), width=110,
                                         text_color="white", text_color_disabled="white",
                                         command=lambda: self.on_button_click("Skip Settings"))
        self.adjust_button.grid(row=row_offset + 1, column=i * 2, padx=(20, 20), pady=5)
        self.skip_button.grid(row=row_offset + 2, column=i * 2, padx=(20, 20), pady=5)
        self.cancel_settings = ctk.CTkButton(self.button_frame, text=self.get_text("cancel"),
                                             font=ctk.CTkFont(size=12), width=110, border_width=1, hover=False,
                                             fg_color="transparent", text_color="#2563EB", border_color="#2563EB",
                                             command=lambda: self.on_button_click("Cancel Settings"))
        self.cancel_settings.grid(row=row_offset + 3, column=i * 2, pady=5, padx=(5, 5))
        self.cancel_settings.grid_remove()
        self.update_button(i, self.adjust_button)
        self.update_button(i, self.skip_button)

    def generation_flow_2(self, i, row_offset=1):
        icon_image, arrow_image = self.update_image(i)

        generate_icon = ctk.CTkButton(self.icon_frame, image=icon_image, text="", width=65,
                                      height=65, fg_color="transparent", hover=False)
        generate_icon.grid(row=0, column=i * 2, padx=0, pady=(50, 0))
        arrow = ctk.CTkButton(self.icon_frame, text='', width=50, height=20, image=arrow_image,
                              fg_color="transparent", hover=False)
        arrow.grid(row=0, column=i * 2 + 1, pady=(50, 0))

        self.progress_bar = ctk.CTkProgressBar(self.button_frame, width=100, height=15, corner_radius=0,
                                               bg_color="transparent", progress_color="#2563EB")
        self.progress_bar.grid(row=row_offset + 1, column=i * 2, padx=5, pady=5)
        self.progress_bar.set(0)

        # Hide progress frame initially
        self.progress_bar.grid_remove()

        self.generate_button = ctk.CTkButton(self.button_frame, text=self.get_text("generate_video"),
                                             font=ctk.CTkFont(size=12), width=110, text_color="white",
                                             text_color_disabled="white",
                                             command=lambda label="Generate Video": self.on_button_click(label))
        self.generate_button.grid(row=row_offset + 1, column=i * 2, pady=5, padx=5)
        self.update_button(i, self.generate_button)

    def review_flow_3(self, i, row_offset=1):
        icon_image, arrow_image = self.update_image(i)
        review_icon = ctk.CTkButton(self.icon_frame, image=icon_image, text="", width=65,
                                    height=65, fg_color="transparent", hover=False)
        review_icon.grid(row=0, column=i * 2, padx=(0, 60), pady=(50, 0))
        self.play_button = ctk.CTkButton(self.button_frame, text=self.get_text("preview_and_play"),
                                         font=ctk.CTkFont(size=12), width=110,
                                         text_color="white", text_color_disabled="white",
                                         command=lambda: self.on_button_click("Preview and Play"))
        self.reselect_button = ctk.CTkButton(self.button_frame, text=self.get_text("reselect_ppt"),
                                             font=ctk.CTkFont(size=12), width=110,
                                             text_color="white", text_color_disabled="white",
                                             command=lambda: self.on_button_click("Reselect PPT"))
        self.play_button.grid(row=row_offset + 1, column=i * 2, pady=5, padx=(20, 40))
        self.reselect_button.grid(row=row_offset + 2, column=i * 2, pady=5, padx=(20, 40))
        self.update_button(i, self.play_button)
        self.update_button(i, self.reselect_button)

    def update_button(self, icon_index, button):
        # 根据步骤更新按钮状态
        if icon_index < self.step:
            # 完成的步骤
            button.configure(state="disabled", fg_color="#B7B7B7", text_color="white")
        elif icon_index == self.step:
            # 当前步骤
            button.configure(state="normal", fg_color="#2563EB", text_color="white", hover_color="#1D4ED8")
        else:
            # 后续步骤，禁用
            button.configure(state="disabled", fg_color="#B7B7B7", text_color="white")

    def update_image(self, icon_index):
        # 根据步骤更新图标
        if icon_index <= self.step:
            # 完成的步骤
            icon_image = self.load_ctk_image(self.icons[icon_index], size=60)
            arrow_image = self.load_ctk_image(self.arrow_icons[0], size=(50, 16))
        else:
            # 后续步骤，禁用
            icon_image = self.load_ctk_image(self.disabled_icons[icon_index], size=60)
            arrow_image = self.load_ctk_image(self.arrow_icons[1], size=(50, 16))
        return icon_image, arrow_image

    def prev_step(self):
        logger.info(f"Current Step: {self.step}")
        if self.step > 0:
            self.step -= 1
            self.update_flow()
        logger.info(f"Go previous Step: {self.step}")

    def next_step(self):
        logger.info(f"Current Step: {self.step}")
        if self.step == 0 and self.file_display == "":
            messagebox.showerror("Error", "Please select a PPT file before proceeding.")
            return
        if self.step < len(self.icons) - 1:
            self.step += 1
            self.update_flow()
        logger.info(f"Go next Step: {self.step}")

    def update_flow(self):
        self.select_flow_0(0)
        self.setting_flow_1(1)
        self.generation_flow_2(2)
        self.review_flow_3(3)

    def center_window(self):
        """Center the window on the screen."""
        screen_width = self.winfo_screenwidth()
        screen_height = self.winfo_screenheight()
        window_width = 800
        window_height = 450
        x = (screen_width - window_width) // 2
        y = (screen_height - window_height) // 2
        self.geometry(f"{window_width}x{window_height}+{x}+{y}")

        self.resizable(False, False)

    def get_text(self, key):
        return self.translations.get(key, key)

    def text_to_key(self, text):
        """Convert display text back to key"""
        for key, value in self.translations.items():
            if value == text:
                return key
        return text

    def change_language(self, language):
        if language in self.language_modes:
            self.current_language = language
            self.translations = self.get_translation()

    def get_translation(self):
        if self.current_language in self._translations_cache:
            return self._translations_cache[self.current_language]

        locale_dir = resource_path(os.path.join('pptflow', os.path.join('locales', self.current_language)))
        translation_file = os.path.join(locale_dir, 'messages.json')

        if not os.path.exists(translation_file):
            locale_dir = resource_path(os.path.join('pptflow', os.path.join('locales', self.language_modes[0])))
            translation_file = os.path.join(locale_dir, 'messages.json')

        with open(translation_file, 'r', encoding='utf-8') as f:
            translations = json.load(f)

        self._translations_cache[self.current_language] = translations
        return translations

    def clear_temp_cache(self):
        self.clear_audio_cache()
        self.clear_image_cache()
        logger.info("Clear temp cache")

    def clear_audio_cache(self):
        if os.path.exists(self.setting.audio_dir_path):
            shutil.rmtree(self.setting.audio_dir_path)
        logger.info("Clear audio cache")

    def clear_image_cache(self):
        if os.path.exists(self.setting.image_dir_path):
            shutil.rmtree(self.setting.image_dir_path)
        logger.info("Clear image cache")

    def on_button_click(self, label_text):
        logger.info(f"Button clicked! Text:{label_text}")
        if label_text == "Select PPT File":
            self.browse_file()
        elif label_text == "Cancel File":
            self.reselect_file()
        elif label_text == "Adjust Settings":
            self.select_frame("Adjust Settings")
        elif label_text == "Skip Settings":
            self.step += 1
            self.adjust_button.grid_remove()
            self.skip_button.grid_remove()
            self.setting_flow_1(1)
            self.cancel_settings.grid()
            self.generation_flow_2(2)
        elif label_text == "Cancel Settings":
            self.step = 1
            self.cancel_settings.grid_remove()
            self.setting_flow_1(1)
            self.generation_flow_2(2)
            self.review_flow_3(3)
        elif label_text == "Generate Video":
            threading.Thread(target=self.generate_video).start()
        elif label_text == "Preview and Play":
            self.play_video()
        elif label_text == "Reselect PPT":
            self.step = 0
            self.file_display = ""
            self.create_workflow_section()
        elif label_text == "Website":
            import webbrowser
            logger.info("Loaded webbrowser")
            webbrowser.open("https://pptflow.com")

    def browse_file(self):
        self.file_display = filedialog.askopenfilename(
            filetypes=[("PowerPoint files", "*.pptx")]
        )
        if self.file_display:
            logger.info(f"Selected file: {self.file_display}")

            # Check whether the PPT file's notes only has English text and punctuation
            if not self.check_ppt_notes_only_english(self.file_display):
                return

            # Set the default output path
            self.setting.video_path = re.sub(r"pptx?$", self.setting.video_format.lower(), self.file_display)

            self.file_label.grid()
            self.file_label.configure(state=ctk.NORMAL)
            # 显示文件名
            display_text = truncate_text(os.path.basename(self.file_display), max_length=10)
            self.tooltip = CustomTooltip(self.file_label, os.path.basename(self.file_display))
            self.file_label_var.set(display_text)
            self.file_label.configure(state=ctk.DISABLED)
            self.cancel_file.grid()

            self.select_button.grid_remove()
            self.step += 1
            self.setting_flow_1(1)
            self.tts = self.tts if self.tts else self.load_tts(self.setting.tts_service_provider)

    def check_ppt_notes_only_english(self, ppt_file_path):
        from pptx import Presentation
        try:
            # ================ 新增校验逻辑 ================
            presentation = Presentation(ppt_file_path)

            # 1. 检查是否有幻灯片
            if len(presentation.slides) == 0:
                messagebox.showerror("Error", "The PPT has no slides, please reselect!")
                logger.error("The PPT has no slides, please reselect!")
                self.reselect_file()
                return False

            # 2. 检查备注和英文内容
            has_notes = False
            non_english_chars = set()

            for slide in presentation.slides:
                notes = slide.notes_slide.notes_text_frame.text.strip()

                # 检查备注是否存在
                if not notes:
                    messagebox.showerror("Error", f"There are slides without notes, please add!")
                    logger.error(f"There are slides without notes, please add!")
                    self.reselect_file()
                    return False

                # 检查是否包含非英文字符
                if not re.match(r'^[\x00-\x7F\u2014\u201C\u201D\u2018\u2019]+$', notes):
                    non_english_chars.update(re.findall(r'[^\x00-\x7F]', notes))

            # 3. 提示非英文字符错误
            if non_english_chars:
                messagebox.showerror("Error",
                                     f"Non-english characters found: {', '.join(non_english_chars)}\n"
                                     "Please change the notes to English only!")
                logger.error(f"Non-english characters found: {', '.join(non_english_chars)}")
                self.reselect_file()
                return False
            return True
        except Exception as e:
            messagebox.showerror("Error", f"Unable to open PPT file: {e}")
            logger.error(f"Unable to open PPT file: {e}", exc_info=True)
            self.reselect_file()
            return False

    def select_frame(self, name):
        self.flow_frame.grid_remove()
        if name == "Adjust Settings":
            from .adjust_settings import AdjustSettingsFrame

            # 显示 ExportSection
            self.adjust_settings = AdjustSettingsFrame(self, self.main_frame)
            self.adjust_settings.grid(row=1, column=0, padx=(100, 40), pady=(40, 60), sticky="nsew")
            self.adjust_settings.tkraise()
            self.adjust_settings.grab_set()
        elif name == "Skip Settings":
            self.next_step()
        elif name == "System Settings":
            if self.adjust_settings is not None:
                self.adjust_settings.grid_remove()
            self.flow_frame.grid_remove()
            from .system_settings import SystemSettingsFrame
            self.system_settings = SystemSettingsFrame(self, self.main_frame)
            self.system_settings.grid(row=1, column=0, padx=(100, 40), pady=(40, 60), sticky="nsew")
            self.system_settings.tkraise()
            self.system_settings.grab_set()

    def load_tts(self, tts_service_provider):
        # import tts module according to service provider
        try:
            tts_service = get_tts_service(tts_service_provider)
            voice_list = tts_service.get_voice_list(self.setting)
            if tts_service_provider == "azure":
                sd.tts_speech_voices = voice_list if len(voice_list) > 0 else sd.tts_speech_voices
            elif tts_service_provider == "kokoro":
                sd.kokoro_voice_type = voice_list if len(voice_list) > 0 else sd.kokoro_voice_type
            return tts_service.tts
        except Exception as e:
            logger.error(f"Error loading TTS service: {e}", exc_info=True)
            messagebox.showerror("Error", f"Error loading TTS service: {e}")
            return None

    def update_progress(self, progress: float, status: str):
        """Update progress bar and status label"""
        self.progress_bar.set(progress)
        # self.status_label.configure(text=status)

    def generate_video(self):
        if not self.file_display:
            messagebox.showerror(self.loading_title, self.get_text("no_file_selected"))
            return
        try:
            # Show progress frame
            self.progress_bar.grid()
            self.generate_button.grid_remove()
            self.update()

            # Load ppt2video
            from pptflow import ppt2video
            logger.info("Loaded ppt2video")
            # Initialize progress tracker
            from pptflow.utils.progress_tracker import ProgressTracker
            logger.info("Loaded ProgressTracker")
            self.progress_tracker = ProgressTracker(self.update_progress)
            ppt2video.ppt_to_video(self.tts, self.file_display, self.setting, self.progress_tracker)

            messagebox.showinfo(self.loading_title,
                                f'{self.get_text("video_generated")}{self.setting.video_path}')
            self.step += 1
            self.generation_flow_2(2)
            self.review_flow_3(3)
        except Exception as e:
            messagebox.showerror("Error", f"Failed to generate video: {str(e)}")
            logger.error(e, exc_info=True)
            return
        finally:
            # Hide progress frame
            self.progress_bar.grid_remove()
            self.generate_button.grid()
            self.progress_bar.set(0)

    def play_video(self):
        logger.info(f'video_path:{self.setting.video_path}')
        try:
            if not os.path.exists(self.setting.video_path):
                logger.error(f'video_path:{self.setting.video_path} does not exist!')
                raise FileNotFoundError(f'video_path:{self.setting.video_path} does not exist!')
            if sys.platform == "win32":
                os.startfile(self.setting.video_path)
            else:
                import subprocess
                logger.info("Loaded subprocess")
                opener = "open" if sys.platform == "darwin" else "xdg-open"
                subprocess.call([opener, self.setting.video_path])
        except Exception as e:
            messagebox.showerror("Error", "No video was generated!")
            logger.error(e, exc_info=True)

    def reselect_file(self):
        self.step = 0
        self.file_label_var.set("")
        self.tooltip = None
        self.setting.video_path = ""
        self.cancel_file.grid_remove()
        self.create_workflow_section()

    def load_ctk_image(self, file_name, size):
        try:
            icon_dir = resource_path(os.path.join("assets", "icons"))
            file_name = os.path.join(icon_dir, file_name)
            # 加载并调整图标大小
            if type(size) == int:
                size = (size, size)
            img = Image.open(file_name)
            img = img.resize(size, Image.LANCZOS)  # 替换 ANTI_ALIAS 为 LANCZOS
            return ctk.CTkImage(light_image=img, dark_image=img, size=size)  # 使用 CTkImage
        except Exception as e:
            print(f"Error loading icon {file_name}: {e}")
            return None

    def update_language(self):
        if self.select_button:
            self.select_button.configure(text=self.get_text("select_ppt"))
        if self.cancel_file:
            self.cancel_file.configure(text=self.get_text("cancel"))
        self.adjust_button.configure(text=self.get_text("adjust_settings"))
        self.skip_button.configure(text=self.get_text("skip_settings"))
        if self.cancel_settings:
            self.cancel_settings.configure(text=self.get_text("cancel"))
        if self.generate_button:
            self.generate_button.configure(text=self.get_text("generate_video"))
        self.play_button.configure(text=self.get_text("preview_and_play"))
        self.reselect_button.configure(text=self.get_text("reselect_ppt"))


def get_locales_subdirectories():
    # Get the absolute path to the current file
    current_file_path = os.path.abspath(__file__)
    # Get the directory path of the current file
    current_directory = os.path.dirname(current_file_path)
    # Get the parent path of the current directory
    parent_directory = os.path.dirname(current_directory)
    # Get the absolute path of the parent directory
    language_locales_path = os.path.join(parent_directory, 'locales')
    try:
        # Get all entries in the directory
        entries = os.listdir(language_locales_path)
        # Filter out subdirectories
        subdirectories = [
            entry for entry in entries
            if os.path.isdir(os.path.join(language_locales_path, entry)) and entry != '__pycache__'
        ]
        return subdirectories
    except FileNotFoundError:
        logger.error(f"Directory {language_locales_path} doesn't exist!")
        return []
    except PermissionError:
        logger.error(f"No permission to access the directory {language_locales_path}")
        return []


def truncate_text(text, max_length=20):
    """ 截取文本，确保头部+省略号+尾部不超过 max_length """
    if len(text) > max_length:
        return text[:max_length // 2] + "..." + text[-(max_length // 2):]
    return text


if __name__ == "__main__":
    app = PPTFlowApp()
    app.mainloop()
